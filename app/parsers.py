from __future__ import annotations

import io
from collections.abc import Iterable


class ParseError(RuntimeError):
    pass


def _clean_text(value: str) -> str:
    return " ".join(value.split())


def _non_empty_units(units: Iterable[dict]) -> list[dict]:
    out: list[dict] = []
    for unit in units:
        text = _clean_text(unit.get("text", ""))
        if text:
            out.append({"text": text, "metadata": unit.get("metadata", {})})
    return out


def parse_pdf(data: bytes) -> list[dict]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ParseError("pypdf is required to parse PDF files") from exc

    reader = PdfReader(io.BytesIO(data))
    units = []
    for i, page in enumerate(reader.pages, start=1):
        units.append(
            {
                "text": page.extract_text() or "",
                "metadata": {"page": i},
            }
        )
    return _non_empty_units(units)


def parse_pptx(data: bytes) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParseError("python-pptx is required to parse PPTX files") from exc

    prs = Presentation(io.BytesIO(data))
    units = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        units.append(
            {
                "text": "\n".join(parts),
                "metadata": {"slide": i},
            }
        )
    return _non_empty_units(units)


def parse_docx(data: bytes) -> list[dict]:
    try:
        from docx import Document
    except ImportError as exc:
        raise ParseError("python-docx is required to parse DOCX files") from exc

    doc = Document(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs if p.text)
    return _non_empty_units([{"text": text, "metadata": {"section": 1}}])


def parse_txt(data: bytes) -> list[dict]:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("latin-1", errors="replace")
    return _non_empty_units([{"text": text, "metadata": {"line_group": 1}}])


def parse_source(data: bytes, source_type: str) -> list[dict]:
    stype = source_type.lower()

    if stype == "pdf":
        units = parse_pdf(data)
    elif stype == "pptx":
        units = parse_pptx(data)
    elif stype == "ppt":
        # Legacy .ppt parsing is out of scope for python-pptx.
        raise ParseError(".ppt is not supported yet. Please convert to .pptx and retry.")
    elif stype == "docx":
        units = parse_docx(data)
    elif stype == "txt":
        units = parse_txt(data)
    elif stype in {"topic", "text"}:
        units = parse_txt(data)
    else:
        raise ParseError(f"Unsupported source type '{source_type}'")

    if not units:
        raise ParseError("No extractable text found in source")

    return units
